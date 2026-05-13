from django.shortcuts import render
from django.http import HttpResponse
from django.template.loader import render_to_string
from django.utils import timezone
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from weasyprint import HTML
from tracker.models import FistulaCase
from mpdsr.models import MPDSREvent
from .ai_utils import generate_newsletter_narrative
from .models import MonthlyNewsletter


def export_pdf(request):
    fistula_operated = FistulaCase.objects.filter(referral_status='OPERATED').count()
    total_deaths = MPDSREvent.objects.count()
    implemented = MPDSREvent.objects.filter(action_status='IMPLEMENTED').count()
    action_gap_percent = (implemented / total_deaths * 100) if total_deaths > 0 else 0

    context = {
        'fistula_operated': fistula_operated,
        'total_deaths': total_deaths,
        'action_gap_percent': action_gap_percent,
    }
    html_string = render_to_string('one_pager.html', context)
    pdf = HTML(string=html_string).write_pdf()
    response = HttpResponse(pdf, content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="report.pdf"'
    return response


def export_ppt(request):
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CIPRB M&E Report"
    slide1.placeholders[1].text = "Executive Summary"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Fistula Campaign Progress"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Surgery Outcomes"
    fistula_operated = FistulaCase.objects.filter(referral_status='OPERATED').count()
    fistula_total = FistulaCase.objects.count()
    p = tf.add_paragraph()
    p.text = f"Total Cases: {fistula_total}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Operated Cases: {fistula_operated}"
    p.level = 1

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "MPDSR Action Tracking"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Data-to-Action Gap"
    total_deaths = MPDSREvent.objects.count()
    implemented = MPDSREvent.objects.filter(action_status='IMPLEMENTED').count()
    action_gap_percent = (implemented / total_deaths * 100) if total_deaths > 0 else 0
    p = tf3.add_paragraph()
    p.text = f"Total MPDSR Events: {total_deaths}"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = f"Implemented Actions: {implemented} ({action_gap_percent:.1f}%)"
    p.level = 1

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    response = HttpResponse(output.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="report.pptx"'
    return response


def generate_newsletter(request):
    now = timezone.now()
    narrative = generate_newsletter_narrative(now.month, now.year)
    newsletter = MonthlyNewsletter.objects.create(month=now.month, year=now.year, content=narrative)
    return render(request, 'newsletter_result.html', {'newsletter': newsletter})
